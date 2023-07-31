from flask import Flask, request, jsonify
from PyPDF2 import PdfReader 
from pptx import Presentation
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from dotenv import load_dotenv
from werkzeug.utils import secure_filename
import os
import pickle
import openai
import requests
from bs4 import BeautifulSoup
from langchain.tools import Tool
from langchain.utilities import GoogleSearchAPIWrapper


openai.api_key =''
os.environ["GOOGLE_CSE_ID"] = ""
os.environ["GOOGLE_API_KEY"] = ""

from flask_cors import CORS

app = Flask(__name__)
CORS(app) 

conversation_chain = None

UPLOAD_FOLDER = r'D:\shree\pdf bot\ask-multiple-pdfs-main\v1\backend' 
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

ALLOWED_EXTENSIONS = {'pdf', 'pptx'}

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_pdf_text(pdf_file):
    pdf_reader = PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

# def scrape_text(url):
#     response = requests.get(url)
#     html_content = response.content
#     soup = BeautifulSoup(html_content, 'html.parser')
#     text = soup.get_text()
    
#     print(text)
#     return text


def scrape_text(url):
    response = requests.get(url)
    html_content = response.content
    soup = BeautifulSoup(html_content, 'html.parser')
    text = soup.get_text()
    
    # with open('output.txt', 'w') as f:
    #     f.write(text)
    
    return text

def get_ppt_text(ppt_file):
    presentation = Presentation(ppt_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
    return text

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings()
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def get_conversation_chain(vectorstore):
    llm = ChatOpenAI(model="gpt-3.5-turbo",max_tokens=1000,temperature=0.4)
    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain

def get_top_products(url):
    company_name = url  
    search = GoogleSearchAPIWrapper(k=10)
    tool = Tool(
        name="Google Search",
        description="Search Google for recent results.",
        func=search.run,
    )
    result = tool.run(f"What is {company_name}'s top products?")
    # print(result)
    return result


@app.route('/upload', methods=['POST'])
def upload():
    global conversation_chain
    load_dotenv()
    if request.content_type == 'application/json':
        data = request.get_json()
        file_type = data['file_type']
        if file_type.lower() == 'url':
            url = data['url']
            raw_text = scrape_text(url)
            
            top_products_text = get_top_products(url)
            raw_text += top_products_text
        else:
            return jsonify({'error': 'Invalid file type. Please enter "url".'}), 400
    elif request.content_type.startswith('multipart/form-data'):
        if 'file' not in request.files:
            return jsonify({'error': 'No file part'}), 400
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No selected file'}), 400
        if file and allowed_file(file.filename):
            if file.content_type == 'application/pdf':
                raw_text = get_pdf_text(file)
            elif file.content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                raw_text = get_ppt_text(file)
            else:
                return jsonify({'error': 'Invalid file type. Please upload a "pdf" or "pptx" file.'}), 400
        else:
            return jsonify({'error': 'File not allowed'}), 400
    else:
        return jsonify({'error': 'Invalid Content-Type.'}), 400


    text_chunks = get_text_chunks(raw_text)
    vectorstore = get_vectorstore(text_chunks)

    # Save the vectorstore to a pkl file
    vectorstore_file_path = os.path.normpath(os.path.join(app.config['UPLOAD_FOLDER'], 'vectorstore.pkl'))
    with open(vectorstore_file_path, 'wb') as f:
        pickle.dump(vectorstore, f)

    conversation_chain = get_conversation_chain(vectorstore)
    return jsonify({'message': 'File uploaded successfully.'}), 200



@app.route('/conversation', methods=['POST'])
def conversation():
    global conversation_chain
    data = request.get_json()
    user_question = data['question']

    # Load the vectorstore from the pkl file
    vectorstore_file_path = os.path.normpath(os.path.join(app.config['UPLOAD_FOLDER'], 'vectorstore.pkl'))
    with open(vectorstore_file_path, 'rb') as f:
     vectorstore = pickle.load(f)


    chat_history = handle_userinput(user_question)
    messages = []
    for i, message in enumerate(chat_history):
        if i % 2 == 0:
            messages.append({'sender': 'user', 'content': message.content})
        else:
            messages.append({'sender': 'bot', 'content': message.content})
    return jsonify({'messages': messages}), 200

def handle_userinput(user_question):
    response = conversation_chain({'question': user_question})
    chat_history = response['chat_history']
    return chat_history



if __name__ == '__main__':
    app.run(port=5000)